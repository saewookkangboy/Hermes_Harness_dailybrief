"""Lecture slide generation — getdesign.md presets, HTML + PPTX."""
from __future__ import annotations

import html
import re
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import yaml

from lib.common import slugify as _slugify

WORKDIR = Path.home() / "hermes-content-studio"


@dataclass
class Slide:
    kind: str  # cover, agenda, body, closing
    title: str
    bullets: list[str]
    notes: str = ""


def load_design_catalog() -> dict[str, Any]:
    path = WORKDIR / "config" / "design-catalog.yaml"
    return yaml.safe_load(path.read_text(encoding="utf-8"))


def select_preset(topic: str, content: str, preset_arg: str | None) -> tuple[str, dict]:
    catalog = load_design_catalog()
    if preset_arg and preset_arg in catalog["presets"]:
        return preset_arg, catalog["presets"][preset_arg]
    text = f"{topic} {content}".lower()
    for rule in catalog.get("selection_rules", []):
        if any(kw.lower() in text for kw in rule.get("keywords", [])):
            name = rule["preset"]
            return name, catalog["presets"][name]
    name = catalog.get("default_preset", "content-studio")
    return name, catalog["presets"][name]


def slugify(s: str) -> str:
    slug = _slugify(s, max_len=40)
    return slug if slug != "weekly" else "lecture"


def parse_content_text(topic: str, raw: str) -> list[Slide]:
    """Parse freeform text into slides."""
    slides: list[Slide] = []
    slides.append(
        Slide(
            kind="cover",
            title=topic,
            bullets=[raw.split("\n")[0][:80] if raw else ""],
            notes="오프닝 — 학습 목표 소개",
        )
    )

    sections = re.split(r"\n(?=#{1,3}\s|\d+\.\s|[─—-]{3,})", raw)
    body_sections: list[tuple[str, list[str]]] = []

    for sec in sections:
        sec = sec.strip()
        if not sec or len(sec) < 10:
            continue
        lines = [ln.strip() for ln in sec.split("\n") if ln.strip()]
        if not lines:
            continue
        title = re.sub(r"^#+\s*", "", lines[0])
        title = re.sub(r"^\d+\.\s*", "", title)
        bullets = []
        for ln in lines[1:]:
            ln = re.sub(r"^[-*•]\s*", "", ln)
            ln = re.sub(r"^\d+\.\s*", "", ln)
            if ln:
                bullets.append(ln[:120])
        if not bullets and len(lines) > 1:
            bullets = [lines[1][:200]]
        if title and title != topic:
            body_sections.append((title[:60], bullets[:5] or [sec[:150]]))

    if not body_sections:
        chunks = [c.strip() for c in re.split(r"\n\n+", raw) if len(c.strip()) > 20]
        for i, chunk in enumerate(chunks[:8], 1):
            first = chunk.split("\n")[0][:60]
            rest = [ln.strip() for ln in chunk.split("\n")[1:4] if ln.strip()]
            body_sections.append((first or f"섹션 {i}", rest or [chunk[:150]]))

    slides.append(
        Slide(
            kind="agenda",
            title="아젠다",
            bullets=[t for t, _ in body_sections[:6]] or ["개념", "사례", "실습", "Q&A"],
            notes="오늘 흐름 안내",
        )
    )

    for title, bullets in body_sections[:10]:
        slides.append(
            Slide(
                kind="body",
                title=title,
                bullets=bullets,
                notes=f"핵심: {bullets[0][:80] if bullets else title}",
            )
        )

    takeaways = body_sections[0][1][:3] if body_sections else ["핵심 개념 정리", "실무 적용", "다음 단계"]
    slides.append(
        Slide(
            kind="closing",
            title="핵심 Takeaway",
            bullets=takeaways,
            notes="마무리 + Q&A",
        )
    )
    return slides


def hex_to_rgb(h: str) -> tuple[int, int, int]:
    h = h.lstrip("#")
    return tuple(int(h[i : i + 2], 16) for i in (0, 2, 4))


def build_html_deck(
    stamp: str,
    topic: str,
    slides: list[Slide],
    preset_name: str,
    preset: dict,
) -> str:
    c = preset["colors"]
    font = preset["typography"]["heading"]
    slide_html = []
    for i, s in enumerate(slides, 1):
        bg = c["tertiary"] if s.kind in ("cover", "closing") else c["surface"]
        fg = c["primary"] if s.kind in ("cover", "closing") else c["primary"]
        if preset_name == "linear" and s.kind in ("cover", "closing"):
            bg = c["neutral"]
            fg = c["primary"]
        band = ""
        if s.kind == "body":
            band = f'<div class="band" style="background:{c["neutral"]}">{html.escape(s.bullets[0][:80] if s.bullets else "")}</div>'
        bullets = ""
        if s.kind != "body" or len(s.bullets) > 1:
            items = "".join(f"<li>{html.escape(b)}</li>" for b in s.bullets)
            bullets = f"<ul>{items}</ul>"
        elif s.kind == "body" and s.bullets:
            bullets = f'<div class="proof">{html.escape(s.bullets[0])}</div>'

        slide_html.append(
            f"""
<section class="slide {s.kind}" style="background:{bg};color:{fg}">
  <div class="bar" style="background:{c['primary']}"></div>
  <h1>{html.escape(s.title)}</h1>
  {band}
  {bullets}
  <span class="page">{i}</span>
  <aside class="notes">{html.escape(s.notes)}</aside>
</section>"""
        )

    return f"""<!DOCTYPE html>
<html lang="ko">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>{html.escape(topic)} — 강의 슬라이드</title>
  <meta name="description" content="{html.escape(topic)}">
  <link rel="stylesheet" href="https://cdn.jsdelivr.net/gh/orioncactus/pretendard/dist/web/static/pretendard.css">
  <style>
    * {{ box-sizing: border-box; margin: 0; padding: 0; }}
    body {{ font-family: '{font}', Pretendard, sans-serif; background: #111; }}
    .deck {{ max-width: 1080px; margin: 0 auto; }}
    .slide {{
      position: relative; min-height: 608px; padding: 48px 56px 48px 72px;
      margin-bottom: 16px; border-radius: 8px; overflow: hidden;
    }}
    .slide .bar {{ position: absolute; left: 0; top: 0; bottom: 0; width: 8px; }}
    .slide.cover .bar, .slide.closing .bar {{ width: 12px; }}
    h1 {{ font-size: 2rem; font-weight: 700; margin-bottom: 1.5rem; letter-spacing: -0.02em; }}
    .band {{ padding: 12px 16px; margin: 0 0 1.5rem; font-size: 0.95rem; border-radius: 4px; }}
    ul {{ margin-left: 1.25rem; line-height: 1.7; }}
    li {{ margin-bottom: 0.5rem; }}
    .proof {{ background: rgba(0,0,0,0.04); padding: 1.5rem; border-radius: 8px; font-size: 1.1rem; }}
    .page {{ position: absolute; bottom: 24px; right: 32px; font-size: 0.75rem; opacity: 0.5; }}
    .notes {{ display: none; }}
    @media print {{ .slide {{ page-break-after: always; margin: 0; border-radius: 0; }} }}
  </style>
</head>
<body>
  <div class="deck" data-preset="{html.escape(preset_name)}" data-design-source="https://getdesign.md/">
    <header style="color:#999;padding:1rem;font-size:0.8rem">
      Design: {html.escape(preset_name)} ({html.escape(preset.get('ref', preset_name))}) · {stamp}
    </header>
    {"".join(slide_html)}
  </div>
</body>
</html>"""


def build_pptx(
    path: Path,
    topic: str,
    slides: list[Slide],
    preset: dict,
) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    c = preset["colors"]

    def add_slide(slide: Slide, idx: int) -> None:
        s = prs.slides.add_slide(blank)
        is_cover = slide.kind in ("cover", "closing")
        bg = c["tertiary"] if is_cover else c["surface"]
        r, g, b = hex_to_rgb(bg)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = RGBColor(r, g, b)

        # Left bar
        bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), prs.slide_height)
        br, bg_c, bb = hex_to_rgb(c["primary"])
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(br, bg_c, bb)
        bar.line.fill.background()

        # Title
        tx = s.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(11.5), Inches(1.2))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = slide.title
        p.font.size = Pt(32 if is_cover else 24)
        p.font.bold = True
        pr, pg, pb = hex_to_rgb(c["primary"])
        p.font.color.rgb = RGBColor(pr, pg, pb)

        # Body band or bullets
        top = 2.0
        if slide.kind == "body" and slide.bullets:
            band = s.shapes.add_shape(1, Inches(0.6), Inches(1.8), Inches(12), Inches(0.6))
            nr, ng, nb = hex_to_rgb(c["neutral"])
            band.fill.solid()
            band.fill.fore_color.rgb = RGBColor(nr, ng, nb)
            band.line.fill.background()
            bt = band.text_frame
            bt.paragraphs[0].text = slide.bullets[0][:100]
            bt.paragraphs[0].font.size = Pt(12)
            top = 2.6
            body_bullets = slide.bullets[1:] or slide.bullets
        else:
            body_bullets = slide.bullets

        bx = s.shapes.add_textbox(Inches(0.6), Inches(top), Inches(11.5), Inches(4))
        btf = bx.text_frame
        btf.word_wrap = True
        for i, bullet in enumerate(body_bullets[:6]):
            para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            para.text = f"• {bullet}"
            para.font.size = Pt(14)
            para.space_after = Pt(8)

        # Page number
        pn = s.shapes.add_textbox(Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.3))
        pn.text_frame.paragraphs[0].text = str(idx)
        pn.text_frame.paragraphs[0].font.size = Pt(10)
        sr, sg, sb = hex_to_rgb(c["secondary"])
        pn.text_frame.paragraphs[0].font.color.rgb = RGBColor(sr, sg, sb)

        # Speaker notes
        if slide.notes:
            s.notes_slide.notes_text_frame.text = slide.notes

    for i, slide in enumerate(slides, 1):
        add_slide(slide, i)

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


def build_outline(stamp: str, topic: str, slides: list[Slide], preset_name: str) -> str:
    lines = [
        f"# 강의 기획 — {topic}",
        f"**날짜:** {stamp}",
        f"**디자인:** {preset_name} ([getdesign.md](https://getdesign.md/))",
        "",
        "## 학습 목표",
        "1. 핵심 개념 이해",
        "2. 실무 적용 방법 습득",
        "3. 다음 단계 액션 플랜 수립",
        "",
        "## 슬라이드 목록",
        "",
        "| # | 유형 | 제목 | 발표자 노트 |",
        "|---|------|------|-------------|",
    ]
    for i, s in enumerate(slides, 1):
        lines.append(f"| {i} | {s.kind} | {s.title} | {s.notes[:40]} |")
    lines.extend(["", "## FAQ (AEO)", "1. ...", "2. ...", "3. ...", ""])
    return "\n".join(lines)


def generate(
    topic: str,
    content: str,
    stamp: str,
    preset_arg: str | None = None,
) -> dict[str, Path]:
    preset_name, preset = select_preset(topic, content, preset_arg)
    slug = slugify(topic)
    out_dir = WORKDIR / "content" / "lectures"
    out_dir.mkdir(parents=True, exist_ok=True)

    slides = parse_content_text(topic, content)
    paths = {
        "outline": out_dir / f"{stamp}_lecture_{slug}_outline.md",
        "html": out_dir / f"{stamp}_lecture_{slug}.html",
        "pptx": out_dir / f"{stamp}_lecture_{slug}.pptx",
    }

    paths["outline"].write_text(
        build_outline(stamp, topic, slides, preset_name), encoding="utf-8"
    )
    paths["html"].write_text(
        build_html_deck(stamp, topic, slides, preset_name, preset), encoding="utf-8"
    )
    build_pptx(paths["pptx"], topic, slides, preset)
    return paths, preset_name
